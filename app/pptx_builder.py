"""
DENSO-branded PPTX builder using python-pptx.
Generates a professional 16:9 presentation with DENSO corporate colors.
"""

from io import BytesIO
from datetime import date
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Logo path — resolve relative to this file
_LOGO_PATH = Path(__file__).parent / "denso_logo.png"

# ---------------------------------------------------------------------------
# DENSO brand palette
# ---------------------------------------------------------------------------
RED = RGBColor(0xC8, 0x10, 0x2E)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x58, 0x59, 0x5B)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

FONT = "Arial"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class DensoPPTXBuilder:
    """Builds a DENSO-branded PowerPoint deck slide by slide."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    # -- reusable elements --------------------------------------------------

    def _blank_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rect(self, slide, left, top, width, height, color):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _top_bar(self, slide, h=Inches(0.8)):
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, h, RED)

    def _bottom_bar(self, slide):
        self._rect(slide, Inches(0), Inches(7.08), SLIDE_W, Inches(0.42), DARK_GRAY)

    def _brand_name(self, slide, left=Inches(0.6), top=Inches(0.15), on_red_bg=False):
        """Add DENSO branding. Uses logo image on white/light backgrounds,
        white text on red backgrounds (logo is red, invisible on red)."""
        if on_red_bg or not _LOGO_PATH.exists():
            tb = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.5))
            p = tb.text_frame.paragraphs[0]
            p.text = "DENSO"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = WHITE if on_red_bg else RED
            p.font.name = FONT
        else:
            # Red logo on white/light background
            slide.shapes.add_picture(
                str(_LOGO_PATH), left, top, Inches(1.6), Inches(0.32),
            )

    def _slide_title(self, slide, text):
        tb = slide.shapes.add_textbox(Inches(3), Inches(0.17), Inches(9.6), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT
        p.alignment = PP_ALIGN.RIGHT

    def _chrome(self, slide, title: str):
        self._top_bar(slide)
        self._brand_name(slide, on_red_bg=True)  # logo sits in red bar → use white text
        self._slide_title(slide, title)
        self._bottom_bar(slide)

    @staticmethod
    def _set_font(paragraph, size=14, bold=False, color=None, italic=False):
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.name = FONT
        paragraph.font.italic = italic
        if color:
            paragraph.font.color.rgb = color

    # -- slide types --------------------------------------------------------

    def add_title_slide(self, title: str, subtitle: str = "", date_str: str | None = None):
        slide = self._blank_slide()
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RED)
        self._brand_name(slide, Inches(0.7), Inches(0.5), on_red_bg=True)

        tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(11.5), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(title.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            self._set_font(p, size=40, bold=True, color=WHITE)

        if subtitle:
            p = tf.add_paragraph()
            p.space_before = Pt(14)
            p.text = subtitle
            self._set_font(p, size=18, color=WHITE)

        self._rect(slide, Inches(0.7), Inches(5.0), Inches(2.2), Inches(0.04), WHITE)

        if date_str:
            tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(5), Inches(0.5))
            p = tb2.text_frame.paragraphs[0]
            p.text = date_str
            self._set_font(p, size=14, color=WHITE)

    def add_section_slide(self, title: str, number: int | None = None):
        slide = self._blank_slide()
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.1), RED)
        self._bottom_bar(slide)
        # Red logo on white background — top right
        self._brand_name(slide, left=Inches(10.8), top=Inches(0.3), on_red_bg=False)

        if number is not None:
            tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(2), Inches(1.2))
            p = tb.text_frame.paragraphs[0]
            p.text = f"{number:02d}"
            self._set_font(p, size=72, bold=True, color=RED)

        tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(11.5), Inches(1.5))
        tf = tb2.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        self._set_font(p, size=34, bold=True, color=BLACK)

    def add_content_slide(self, title: str, body: str = "", bullets: list[str] | None = None):
        slide = self._blank_slide()
        self._chrome(slide, title)

        tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(11.9), Inches(5.6))
        tf = tb.text_frame
        tf.word_wrap = True

        if body:
            for i, para in enumerate(body.split("\n\n")):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = para
                self._set_font(p, size=15, color=DARK_GRAY)
                p.space_after = Pt(10)

        if bullets:
            for b in bullets:
                p = tf.add_paragraph()
                p.text = f"\u25b8  {b}"
                self._set_font(p, size=14, color=DARK_GRAY)
                p.space_before = Pt(5)

    def add_image_slide(self, title: str, image_source, caption: str = ""):
        """image_source can be a file path (str) or BytesIO."""
        slide = self._blank_slide()
        self._chrome(slide, title)

        slide.shapes.add_picture(
            image_source, Inches(1.2), Inches(1.05), Inches(10.9), Inches(5.3),
        )

        if caption:
            tb = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5))
            p = tb.text_frame.paragraphs[0]
            p.text = caption
            self._set_font(p, size=10, italic=True, color=MED_GRAY)
            p.alignment = PP_ALIGN.CENTER

    def add_two_col_slide(
        self, title: str, left_text: str = "",
        left_bullets: list[str] | None = None, right_image=None,
    ):
        slide = self._blank_slide()
        self._chrome(slide, title)

        tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(5.6), Inches(5.6))
        tf = tb.text_frame
        tf.word_wrap = True

        if left_text:
            p = tf.paragraphs[0]
            p.text = left_text
            self._set_font(p, size=14, color=DARK_GRAY)
            p.space_after = Pt(8)

        if left_bullets:
            for b in left_bullets:
                p = tf.add_paragraph()
                p.text = f"\u25b8  {b}"
                self._set_font(p, size=13, color=DARK_GRAY)
                p.space_before = Pt(4)

        if right_image:
            slide.shapes.add_picture(
                right_image, Inches(6.7), Inches(1.15), Inches(6.0), Inches(5.3),
            )

    def add_table_slide(self, title: str, df: pd.DataFrame):
        slide = self._blank_slide()
        self._chrome(slide, title)

        rows, cols = df.shape
        rows += 1
        col_w = min(2.0, 11.9 / cols)
        table_w = Inches(col_w * cols)
        table_h = Inches(min(5.5, rows * 0.45))
        left = (SLIDE_W - table_w) // 2

        tbl = slide.shapes.add_table(
            rows, cols, left, Inches(1.15), table_w, table_h,
        ).table

        for j, col_name in enumerate(df.columns):
            cell = tbl.cell(0, j)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RED
            for para in cell.text_frame.paragraphs:
                self._set_font(para, size=11, bold=True, color=WHITE)
                para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for i in range(df.shape[0]):
            for j in range(df.shape[1]):
                cell = tbl.cell(i + 1, j)
                cell.text = str(df.iloc[i, j])
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE
                for para in cell.text_frame.paragraphs:
                    self._set_font(para, size=10, color=DARK_GRAY)
                    para.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def add_closing_slide(self, title: str = "Thank You", lines: list[str] | None = None):
        slide = self._blank_slide()
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RED)
        self._brand_name(slide, Inches(0.7), Inches(0.5), on_red_bg=True)

        tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(11.5), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        self._set_font(p, size=44, bold=True, color=WHITE)

        self._rect(slide, Inches(0.7), Inches(4.2), Inches(2.2), Inches(0.04), WHITE)

        if lines:
            tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(11.5), Inches(2.5))
            tf = tb2.text_frame
            for ln in lines:
                p = tf.add_paragraph()
                p.text = ln
                self._set_font(p, size=14, color=WHITE)
                p.space_before = Pt(4)

    # -- output -------------------------------------------------------------

    def save_bytes(self) -> BytesIO:
        buf = BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        return buf

    def save_file(self, path: str):
        self.prs.save(path)
